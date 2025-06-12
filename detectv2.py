import torch
from torch import nn
import torch.nn.functional as F
from transformers import BertModel, BertTokenizer, ViTModel, ViTImageProcessor
import os
import sys
import numpy as np
from PIL import Image
import fitz  # PyMuPDF
import io
import re
import docx

# Force offline mode
os.environ["TRANSFORMERS_OFFLINE"] = "1"

# Paths
VIT_MODEL_PATH = "./local_vit_model"
BERT_MODEL_PATH = "./local_bert_model"
MODEL_PATH = "best_model.pt"


# Copy exact model class from train.py
class MultimodalCVClassifier(nn.Module):
    def __init__(self, vit_model_name, bert_model_name, num_classes, num_jobs, dropout_rate=0.3):
        super(MultimodalCVClassifier, self).__init__()
        
        self.vit = ViTModel.from_pretrained(vit_model_name)
        self.bert = BertModel.from_pretrained(bert_model_name)
        
        self.vit_dim = self.vit.config.hidden_size
        self.bert_dim = self.bert.config.hidden_size
        self.format_dim = 4
        self.job_dim = num_jobs
        
        self.visual_branch = nn.Sequential(
            nn.Linear(self.vit_dim, 512),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(512, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        self.visual_score = nn.Linear(256, 1)
        
        self.semantic_branch = nn.Sequential(
            nn.Linear(self.bert_dim * 2, 768),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(768, 512),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(512, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        self.cv_jd_attention = nn.MultiheadAttention(
            embed_dim=self.bert_dim,
            num_heads=8,
            dropout=dropout_rate
        )
        
        self.semantic_score = nn.Linear(256, 1)
        
        self.format_branch = nn.Sequential(
            nn.Linear(self.format_dim, 32),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(32, 64),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        self.job_branch = nn.Sequential(
            nn.Linear(self.job_dim, 32),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(32, 64),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        self.combined_dim = 256 + 256 + 64 + 64
        
        self.attention = nn.Sequential(
            nn.Linear(self.combined_dim, 4),
            nn.Softmax(dim=1)
        )
        
        self.classifier = nn.Sequential(
            nn.Linear(self.combined_dim, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(256, num_classes),
        )

    def forward(self, pixel_values, input_ids, attention_mask, token_type_ids=None, 
                format_features=None, job_one_hot=None):
        
        if format_features is None:
            format_features = torch.zeros(pixel_values.size(0), self.format_dim, 
                                         device=pixel_values.device)
        
        if job_one_hot is None:
            job_one_hot = torch.zeros(pixel_values.size(0), self.job_dim,
                                     device=pixel_values.device)
        
        vit_outputs = self.vit(pixel_values=pixel_values)
        vit_features = vit_outputs.pooler_output
        visual_features = self.visual_branch(vit_features)
        visual_score_val = self.visual_score(visual_features)
        
        bert_outputs = self.bert(
            input_ids=input_ids,
            attention_mask=attention_mask,
            token_type_ids=token_type_ids,
        )
        
        all_hidden_states = bert_outputs.last_hidden_state
        
        jd_mask = (token_type_ids == 0).float()
        cv_mask = (token_type_ids == 1).float()
        
        jd_embeddings = (all_hidden_states * jd_mask.unsqueeze(-1)).sum(dim=1) / jd_mask.sum(dim=1, keepdim=True).clamp(min=1)
        cv_embeddings = (all_hidden_states * cv_mask.unsqueeze(-1)).sum(dim=1) / cv_mask.sum(dim=1, keepdim=True).clamp(min=1)
        
        cv_attended, _ = self.cv_jd_attention(
            cv_embeddings.unsqueeze(0),
            jd_embeddings.unsqueeze(0),
            jd_embeddings.unsqueeze(0),
        )
        cv_attended = cv_attended.squeeze(0)
        
        combined_semantic = torch.cat([cv_embeddings, cv_attended], dim=1)
        semantic_features = self.semantic_branch(combined_semantic)
        semantic_score_val = self.semantic_score(semantic_features)
        
        format_features = self.format_branch(format_features)
        job_features = self.job_branch(job_one_hot)
        
        combined_features = torch.cat(
            (visual_features, semantic_features, format_features, job_features), dim=1
        )
        
        attention_weights = self.attention(combined_features)
        
        semantic_importance = torch.sigmoid(semantic_score_val).squeeze()
        visual_penalty = torch.sigmoid(-visual_score_val).squeeze() * 0.5
        
        adjusted_weights = attention_weights.clone()
        adjusted_weights[:, 1] = adjusted_weights[:, 1] + semantic_importance * 0.3
        adjusted_weights[:, 0] = adjusted_weights[:, 0] - visual_penalty * 0.2
        
        attention_weights = F.softmax(adjusted_weights, dim=1)
        
        visual_weight = attention_weights[:, 0].unsqueeze(1).expand_as(visual_features)
        semantic_weight = attention_weights[:, 1].unsqueeze(1).expand_as(semantic_features)
        format_weight = attention_weights[:, 2].unsqueeze(1).expand_as(format_features)
        job_weight = attention_weights[:, 3].unsqueeze(1).expand_as(job_features)
        
        weighted_visual = visual_features * visual_weight
        weighted_semantic = semantic_features * semantic_weight
        weighted_format = format_features * format_weight
        weighted_job = job_features * job_weight
        
        weighted_features = torch.cat(
            (weighted_visual, weighted_semantic, weighted_format, weighted_job), dim=1
        )
        
        logits = self.classifier(weighted_features)
        
        return {
            'logits': logits,
            'visual_score': visual_score_val,
            'semantic_score': semantic_score_val,
            'attention_weights': attention_weights
        }


def create_page_image(file_path, page_num=0, scale=1.5):
    """Create image from document - same as train.py"""
    if file_path.lower().endswith('.pdf'):
        try:
            doc = fitz.open(file_path)
            if len(doc) > 0:
                page = doc[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                img_data = pix.pil_tobytes("RGB")
                img = Image.open(io.BytesIO(img_data))
                doc.close()
                return img
            doc.close()
        except:
            pass
    
    return Image.new('RGB', (224, 224), color='white')


def extract_text_from_file(file_path):
    """Extract text - same as train.py"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if not os.path.exists(file_path):
            return f"File not found: {os.path.basename(file_path)}"
        
        if file_ext == '.pdf':
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            return text if text.strip() else "PDF content could not be extracted"
            
        elif file_ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
                
        elif file_ext in ['.docx', '.doc']:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
            
        elif file_ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
            
        elif file_ext == '.odt':
            from odf import text, teletype
            from odf.opendocument import load
            doc = load(file_path)
            allparas = doc.getElementsByType(text.P)
            return "\n".join([teletype.extractText(p) for p in allparas])
            
        elif file_ext == '.eml':
            import email
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                msg = email.message_from_file(f)
                body = ""
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body += part.get_payload(decode=True).decode('utf-8', errors='replace')
                return body if body else "No text content in email"
                
        elif file_ext == '.mht':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
                text = re.sub('<.*?>', '', content)
                return text
                
        elif file_ext == '.heic':
            return "Image file (no text)"
            
        elif file_ext == '.bin':
            return "Binary file (no text)"
            
        else:
            return f"Unsupported file type {file_ext}"
            
    except Exception as e:
        return f"Error reading file: {os.path.basename(file_path)}"


def extract_cv_features(file_path):
    """Extract format features - same as train.py"""
    if file_path.lower().endswith('.pdf'):
        try:
            doc = fitz.open(file_path)
            pages = len(doc)
            
            # Simple feature extraction
            text = ""
            for page in doc:
                text += page.get_text()
            
            doc.close()
            
            # Calculate features
            bullet_points = text.count('•') + text.count('▪') + text.count('-')
            text_density = len(text) / max(pages, 1) / 1000
            
            return {
                "font_count": 2,  # Default
                "text_density": text_density,
                "bullet_points": bullet_points,
                "pages": pages
            }
        except:
            pass
    
    # Default values for non-PDFs
    return {
        "font_count": 1,
        "text_density": 1,
        "bullet_points": 0,
        "pages": 1
    }


def predict_cv_status(jd_path, cv_path):
    """Main prediction function"""
    
    # Device setup
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")
    
    # Load tokenizer and processor
    print("Loading models...")
    bert_tokenizer = BertTokenizer.from_pretrained(BERT_MODEL_PATH)
    vit_processor = ViTImageProcessor.from_pretrained(VIT_MODEL_PATH)
    
    # Initialize model (10 jobs as per training)
    model = MultimodalCVClassifier(
        vit_model_name=VIT_MODEL_PATH,
        bert_model_name=BERT_MODEL_PATH,
        num_classes=4,
        num_jobs=10
    )
    
    # Load trained weights
    if os.path.exists(MODEL_PATH):
        model.load_state_dict(torch.load(MODEL_PATH, map_location=device))
        print(f"Model loaded from {MODEL_PATH}")
    else:
        print(f"ERROR: Model file {MODEL_PATH} not found!")
        return
    
    model.to(device)
    model.eval()
    
    # Extract texts
    print(f"\nReading JD: {jd_path}")
    jd_text = extract_text_from_file(jd_path)
    
    print(f"Reading CV: {cv_path}")
    cv_text = extract_text_from_file(cv_path)
    
    # Process text with BERT
    encoding = bert_tokenizer(
        jd_text,
        cv_text,
        add_special_tokens=True,
        max_length=512,  # BERT max length
        padding="max_length",
        truncation=True,
        return_tensors="pt",
    )
    
    # Process image with ViT
    cv_image = create_page_image(cv_path)
    visual_inputs = vit_processor(images=cv_image, return_tensors="pt")
    
    # Extract format features
    features = extract_cv_features(cv_path)
    format_features = torch.tensor([
        min(features["font_count"], 10) / 10.0,
        min(features["text_density"], 20) / 20.0,
        min(features["bullet_points"], 30) / 30.0,
        min(features["pages"], 5) / 5.0
    ], dtype=torch.float32).unsqueeze(0)
    
    # Move to device
    input_ids = encoding["input_ids"].to(device)
    attention_mask = encoding["attention_mask"].to(device)
    token_type_ids = encoding["token_type_ids"].to(device)
    pixel_values = visual_inputs["pixel_values"].to(device)
    format_features = format_features.to(device)
    
    # Predict
    print("\nPredicting...")
    with torch.no_grad():
        outputs = model(
            pixel_values=pixel_values,
            input_ids=input_ids,
            attention_mask=attention_mask,
            token_type_ids=token_type_ids,
            format_features=format_features
        )
        
        logits = outputs['logits']
        visual_score = outputs['visual_score'].item()
        semantic_score = outputs['semantic_score'].item()
        attention_weights = outputs['attention_weights'][0]
        
        # Get prediction
        probabilities = F.softmax(logits, dim=1)
        predicted_class = torch.argmax(probabilities, dim=1).item()
        confidence = probabilities[0][predicted_class].item()
    
    # Status mapping
    status_map = {0: "ACCEPT", 1: "INTERVIEW", 2: "SHORTLIST", 3: "REJECT"}
    
    # Display results
    print("\n" + "="*60)
    print("CV EVALUATION RESULTS")
    print("="*60)
    print(f"Status: {status_map[predicted_class]}")
    print(f"Confidence: {confidence:.2%}")
    print(f"\nScores:")
    print(f"  Visual Score: {1/(1+np.exp(-visual_score)):.2%}")
    print(f"  Semantic Score: {1/(1+np.exp(-semantic_score)):.2%}")
    print(f"\nProbabilities:")
    for i, status in status_map.items():
        print(f"  {status}: {probabilities[0][i].item():.2%}")
    print("="*60)


def main():
    if len(sys.argv) != 3:
        print("Usage: python detectv2.py <job_description_file> <cv_file>")
        print("Example: python detectv2.py job.txt resume.pdf")
        sys.exit(1)
    
    jd_path = sys.argv[1]
    cv_path = sys.argv[2]
    
    if not os.path.exists(jd_path):
        print(f"Error: JD file not found: {jd_path}")
        sys.exit(1)
        
    if not os.path.exists(cv_path):
        print(f"Error: CV file not found: {cv_path}")
        sys.exit(1)
    
    predict_cv_status(jd_path, cv_path)


if __name__ == "__main__":
    main()