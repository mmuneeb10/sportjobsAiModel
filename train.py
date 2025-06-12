import torch
from torch import nn
import torch.optim as optim
import torch.nn.functional as F
from torch.utils.data import Dataset, DataLoader, random_split
from transformers import ViTImageProcessor, ViTModel, BertModel, BertTokenizer
import numpy as np
import os
import json
import fitz  # PyMuPDF for PDF processing
from PIL import Image
import io
import random
import re
import docx  # For reading Word documents
import warnings
import logging

# Suppress the overflowing tokens warning
logging.getLogger("transformers.tokenization_utils_base").setLevel(logging.ERROR)

# Force offline mode to use local models
os.environ["TRANSFORMERS_OFFLINE"] = "1"

# Paths to local models
VIT_MODEL_PATH = "./local_vit_model"
BERT_MODEL_PATH = "./local_bert_model"
CV_FILES_DIR = "./updatedCvs"

# Define status/label mapping 
def map_status_to_label(status):
    """
    Maps job application status to numeric class label
    
    Args:
        status (str): Status string like "ACCEPT", "REJECT", etc.
        
    Returns:
        int: Numeric class label (0-4)
    """
    status_map = {
        "ACCEPT": 0,
        "INTERVIEW": 1, 
        "SHORTLIST": 2,
        "REJECT": 3,
    }
    # Convert to uppercase and get value (default to 0 if not found)
    return status_map.get(status.upper(), 0)  

def create_page_image(file_path, page_num=0, scale=1.5):
    """
    Extracts a specific page from a document and converts it to a PIL Image
    
    Supports PDF files. For other formats, will create a blank image.
    
    Args:
        file_path (str): Path to document file
        page_num (int): Page number to extract (0-indexed)
        scale (float): Scale factor for rendering (higher = better quality)
        
    Returns:
        PIL.Image: Image of the rendered page
    """
    # Handle PDF files
    if file_path.lower().endswith('.pdf'):
        doc = None
        try:
            # Check if file exists first
            if not os.path.exists(file_path):
                return Image.new('RGB', (224, 224), color='white')
            
            # Check if file is readable
            if not os.access(file_path, os.R_OK):
                return Image.new('RGB', (224, 224), color='white')
            
            # Try to open the PDF with PyMuPDF
            try:
                doc = fitz.open(file_path)
            except fitz.FileDataError as fde:
                return Image.new('RGB', (224, 224), color='white')
            except fitz.FileNotFoundError as fnf:
                return Image.new('RGB', (224, 224), color='white')
            except Exception as open_error:
                return Image.new('RGB', (224, 224), color='white')
            
            # Check if document has pages
            if len(doc) == 0:
                doc.close()
                return Image.new('RGB', (224, 224), color='white')
            
            # Check if page exists
            if page_num >= len(doc):
                page_num = 0  # Default to first page
                
            # Get the page
            page = doc[page_num]
            
            # Render page to an image (RGB) with error handling
            try:
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            except Exception as render_error:
                doc.close()
                return Image.new('RGB', (224, 224), color='white')
            
            # Convert to PIL Image
            try:
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            except Exception as img_error:
                doc.close()
                return Image.new('RGB', (224, 224), color='white')
            
            # Resize to fit ViT input requirements (224x224)
            # Use LANCZOS for backward compatibility with older PIL versions
            try:
                img = img.resize((224, 224), Image.Resampling.LANCZOS)
            except AttributeError:
                # Fallback for older PIL versions
                img = img.resize((224, 224), Image.LANCZOS)
            
            doc.close()
            return img
        
        except Exception as e:
            # Make sure to close the document if it was opened
            if doc is not None:
                try:
                    doc.close()
                except:
                    pass
            # Return a blank image in case of error
            return Image.new('RGB', (224, 224), color='white')
    else:
        # For non-PDF files, create a blank image
        # In a production setting, you might want to add support for other formats
        return Image.new('RGB', (224, 224), color='white')

def extract_text_from_file(file_path):
    """
    Extract text content from various file types (PDF, TXT, DOCX)
    
    Args:
        file_path (str): Path to file
        
    Returns:
        str: Extracted text content
    """
    # Check file extension
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        # Check if file exists
        if not os.path.exists(file_path):
            return f"File not found: {os.path.basename(file_path)}"
        
        # Check if file is readable
        if not os.access(file_path, os.R_OK):
            return f"File not readable: {os.path.basename(file_path)}"
        
        # PDF files
        if file_ext == '.pdf':
            doc = None
            try:
                doc = fitz.open(file_path)
                text = ""
                # Extract text from all pages
                for page_num in range(len(doc)):
                    try:
                        page = doc[page_num]
                        page_text = page.get_text()
                        text += page_text + "\n"
                    except Exception as page_error:
                        continue
                
                # If no text was extracted, try alternative extraction methods
                if not text.strip():
                    # Try extracting text blocks
                    for page_num in range(len(doc)):
                        try:
                            page = doc[page_num]
                            blocks = page.get_text("blocks")
                            for block in blocks:
                                if block[6] == 0:  # Text block
                                    text += block[4] + "\n"
                        except:
                            pass
                
                doc.close()
                
                # If still no text was extracted, return a placeholder
                if not text.strip():
                    return f"PDF content could not be extracted from {os.path.basename(file_path)}"
                
                return text
            except fitz.FileDataError as fde:
                return f"Corrupted PDF: {os.path.basename(file_path)}"
            except fitz.FileNotFoundError as fnf:
                return f"PDF not found: {os.path.basename(file_path)}"
            except Exception as pdf_error:
                return f"Failed to read PDF: {os.path.basename(file_path)}"
            finally:
                # Ensure document is closed even if there's an error
                if doc is not None:
                    try:
                        doc.close()
                    except:
                        pass
            
        # Text files
        elif file_ext == '.txt':
            try:
                # Try multiple encodings
                encodings = ['utf-8', 'latin-1', 'windows-1252', 'ascii']
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            return f.read()
                    except UnicodeDecodeError:
                        continue
                # If all encodings fail, use replace errors
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
            except Exception as txt_error:
                return f"Failed to read text file: {os.path.basename(file_path)}"
                
        # Word documents
        elif file_ext in ['.docx', '.doc']:
            try:
                doc = docx.Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            except Exception as doc_error:
                return f"Failed to read Word document: {os.path.basename(file_path)}"
        
        # PowerPoint presentations
        elif file_ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            except:
                return f"Failed to read PowerPoint: {os.path.basename(file_path)}"
        
        # OpenDocument Text
        elif file_ext == '.odt':
            try:
                from odf import text, teletype
                from odf.opendocument import load
                doc = load(file_path)
                allparas = doc.getElementsByType(text.P)
                return "\n".join([teletype.extractText(p) for p in allparas])
            except:
                return f"Failed to read ODT: {os.path.basename(file_path)}"
        
        # Email files
        elif file_ext == '.eml':
            try:
                import email
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    msg = email.message_from_file(f)
                    # Extract body
                    body = ""
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            body += part.get_payload(decode=True).decode('utf-8', errors='replace')
                    return body if body else f"No text content in email: {os.path.basename(file_path)}"
            except:
                return f"Failed to read email: {os.path.basename(file_path)}"
        
        # MHTML files
        elif file_ext == '.mht':
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    content = f.read()
                    # Basic extraction - remove HTML tags
                    text = re.sub('<.*?>', '', content)
                    return text
            except:
                return f"Failed to read MHT: {os.path.basename(file_path)}"
        
        # HEIC images - skip text extraction
        elif file_ext == '.heic':
            return f"Image file (no text): {os.path.basename(file_path)}"
        
        # Binary files
        elif file_ext == '.bin':
            return f"Binary file (no text): {os.path.basename(file_path)}"
        
        # Other file types - return placeholder
        else:
            return f"Unsupported file type {file_ext}: {os.path.basename(file_path)}"
            
    except Exception as e:
        return f"Error reading file: {os.path.basename(file_path)}"

def extract_cv_features(file_path):
    """
    Extract formatting features from a CV to provide additional inputs
    
    For PDFs:
    - Font count
    - Text density
    - Use of bullet points
    - Page count
    
    For other files:
    - Default values
    
    Args:
        file_path (str): Path to CV file
        
    Returns:
        dict: Dictionary of extracted features
    """
    # Only extract detailed features from PDFs
    if file_path.lower().endswith('.pdf'):
        doc = None
        try:
            # Check if file exists
            if not os.path.exists(file_path):
                return {
                    "font_count": 1,
                    "text_density": 1,
                    "bullet_points": 0,
                    "pages": 1
                }
            
            # Check if file is readable
            if not os.access(file_path, os.R_OK):
                return {
                    "font_count": 1,
                    "text_density": 1,
                    "bullet_points": 0,
                    "pages": 1
                }
            
            try:
                doc = fitz.open(file_path)
            except fitz.FileDataError as fde:
                return {
                    "font_count": 1,
                    "text_density": 1,
                    "bullet_points": 0,
                    "pages": 1
                }
            except Exception as open_error:
                return {
                    "font_count": 1,
                    "text_density": 1,
                    "bullet_points": 0,
                    "pages": 1
                }
            
            # Check if document has pages
            if len(doc) == 0:
                doc.close()
                return {
                    "font_count": 1,
                    "text_density": 1,
                    "bullet_points": 0,
                    "pages": 1
                }
            
            # Initialize counters and metrics
            font_counter = {}
            text_blocks = 0
            bullet_points = 0
            page_count = len(doc)
            
            # Process each page
            for page_num in range(len(doc)):
                try:
                    page = doc[page_num]
                    
                    # Extract text blocks
                    try:
                        blocks = page.get_text("dict")["blocks"]
                        text_blocks += len([b for b in blocks if b["type"] == 0])  # Text blocks
                        
                        # Count different fonts
                        for block in blocks:
                            if block["type"] == 0:  # Text block
                                for line in block.get("lines", []):
                                    for span in line.get("spans", []):
                                        font = span.get("font", "unknown")
                                        if font not in font_counter:
                                            font_counter[font] = 0
                                        font_counter[font] += 1
                    except Exception as block_error:
                        pass
                    
                    # Count bullet points (rough estimate)
                    try:
                        text = page.get_text()
                        bullet_points += text.count("•") + text.count("-") + text.count("*")
                    except Exception as text_error:
                        pass
                        
                except Exception as page_error:
                    continue
                
            # Calculate metrics
            font_count = len(font_counter) if font_counter else 1
            text_density = text_blocks / max(1, page_count)  # Blocks per page
            
            doc.close()
            
            return {
                "font_count": max(1, font_count),  # At least 1 font
                "text_density": max(0.1, text_density),  # At least some density
                "bullet_points": max(0, bullet_points),  # Non-negative
                "pages": page_count
            }
            
        except Exception as e:
            # Ensure document is closed even if there's an error
            if doc is not None:
                try:
                    doc.close()
                except:
                    pass
            return {
                "font_count": 1,
                "text_density": 1,
                "bullet_points": 0,
                "pages": 1
            }
    else:
        # Default values for non-PDF files
        return {
            "font_count": 1,  # Assume 1 font
            "text_density": 1, # Medium density
            "bullet_points": 0, # No bullet points 
            "pages": 1  # Single page
        }

def load_cv_data_from_folder_structure(base_dir):
    """
    Load CV data from the specified folder structure.
    
    Expected structure:
    - CV (base_dir)
      - JOB1 (job folder)
        - job_description.docx (or any job description file)
        - ACCEPT (status folder)
          - cv1.pdf, cv2.docx, etc. (CV files)
        - REJECT (status folder)
          - cv3.pdf, etc.
      - JOB2
        - job_description.docx
        - ...and so on
    
    Args:
        base_dir (str): Path to the base CV directory
        
    Returns:
        list: List of (cv_path, job_name, status, label) tuples
    """
    cv_data = []
    
    # Check if base directory exists
    if not os.path.exists(base_dir):
        raise ValueError(f"Base directory {base_dir} does not exist")
    
    # Statistics for tracking
    skipped_files = 0
    processed_jobs = 0
    
    # Iterate through job folders
    for job_dir in sorted(os.listdir(base_dir)):
        job_path = os.path.join(base_dir, job_dir)
        
        # Skip if not a directory or hidden
        if not os.path.isdir(job_path) or job_dir.startswith('.'):
            continue
        
        processed_jobs += 1
        
        # Find job description file - look for a file named jobDescription or job_description
        job_description_file = None
        try:
            for file in os.listdir(job_path):
                file_path = os.path.join(job_path, file)
                if os.path.isfile(file_path) and not file.startswith('.') and \
                   (file.lower() == 'jobdescription' or 'job_description' in file.lower() or \
                    'jobdescription' in file.lower()):
                    job_description_file = file_path
                    break
                    
            # If no specific job description file found, try to find any .txt or .docx file
            if job_description_file is None:
                for file in os.listdir(job_path):
                    file_path = os.path.join(job_path, file)
                    if os.path.isfile(file_path) and not file.startswith('.') and \
                       (file.endswith('.txt') or file.endswith('.docx')):
                        job_description_file = file_path
                        break
        except Exception as e:
            continue
                
        # Valid status folders
        valid_statuses = ['ACCEPT', 'INTERVIEW', 'SHORTLIST', 'REJECT']
        
        # Iterate through status folders
        try:
            for status_dir in os.listdir(job_path):
                status_path = os.path.join(job_path, status_dir)
                
                # Skip if not a directory or if it's a hidden directory
                if not os.path.isdir(status_path) or status_dir.startswith('.'):
                    continue
                
                # Check if it's a valid status folder
                if status_dir.upper() not in valid_statuses:
                    continue
                
                # Map status folder name to label
                label = map_status_to_label(status_dir)
                
                # Iterate through CV files in the status folder
                try:
                    for cv_file in sorted(os.listdir(status_path)):
                        cv_path = os.path.join(status_path, cv_file)
                        
                        # Skip directories and hidden files
                        if os.path.isdir(cv_path) or cv_file.startswith('.'):
                            continue
                        
                        # Check if it's a supported file type
                        file_ext = os.path.splitext(cv_file)[1].lower()
                        supported_exts = ['.pdf', '.txt', '.docx', '.doc']
                        if file_ext not in supported_exts:
                            skipped_files += 1
                            continue
                        
                        # Verify file is readable
                        if not os.access(cv_path, os.R_OK):
                            skipped_files += 1
                            continue
                        
                        # Add to data list with job description file
                        cv_data.append((cv_path, job_dir, status_dir, label, job_description_file))
                except Exception as e:
                    continue
        except Exception as e:
            continue
    
    print(f"\nProcessed {processed_jobs} job folders")
    print(f"Found {len(cv_data)} CV files")
    print(f"Skipped {skipped_files} unsupported/unreadable files")
    
    if not cv_data:
        raise ValueError(f"No CV data found in {base_dir}")
    
    return cv_data

class CVJobDataset(Dataset):
    """
    Dataset for CV data from our folder structure.
    
    Extracts:
    1. Visual features (from PDF or blank image for non-PDFs)
    2. Text features (CV text)
    3. Format features (additional PDF features)
    4. Job information (one hot encoded job name)
    """
    def __init__(self, cv_data, vit_processor, bert_tokenizer, max_length=768, extract_format_features=True, use_chunk_processing=True):
        """
        Initialize the dataset
        
        Args:
            cv_data (list): List of (cv_path, job_name, status, label, job_description_file) tuples
            vit_processor: ViT image processor
            bert_tokenizer: BERT tokenizer
            max_length (int): Max token length for BERT (default 512)
            extract_format_features (bool): Whether to extract additional CV formatting features
            use_chunk_processing (bool): Whether to use chunk processing for long texts
        """
        self.cv_data = cv_data
        self.vit_processor = vit_processor
        self.bert_tokenizer = bert_tokenizer
        self.max_length = max_length
        self.extract_format_features = extract_format_features
        self.use_chunk_processing = use_chunk_processing
        
        # Cache for extracted features and text
        self.format_features_cache = {}
        self.cv_text_cache = {}
        
        # Get unique job names for one-hot encoding
        self.unique_jobs = sorted(set(job for _, job, _, _, _ in cv_data))
        self.job_to_idx = {job: idx for idx, job in enumerate(self.unique_jobs)}
        
        # Cache for job descriptions
        self.job_description_cache = {}
        
    def __len__(self):
        return len(self.cv_data)

    def __getitem__(self, idx):
        cv_path, job_name, status, label, job_description_file = self.cv_data[idx]
        
        # Extract CV text (with caching)
        if cv_path in self.cv_text_cache:
            cv_text = self.cv_text_cache[cv_path]
        else:
            cv_text = extract_text_from_file(cv_path)
            # Ensure we have some text content
            if not cv_text or len(cv_text.strip()) < 10:
                cv_text = f"CV content from {os.path.basename(cv_path)}"
            self.cv_text_cache[cv_path] = cv_text
        
        # Extract job description from saved file (with caching)
        # Instead of generating a synthetic job description
        if job_description_file in self.job_description_cache:
            job_description = self.job_description_cache[job_description_file]
        else:
            # Read the actual job description from the file
            if job_description_file and os.path.exists(job_description_file):
                job_description = extract_text_from_file(job_description_file)
                # If the job description is empty or couldn't be read, use a fallback
                if not job_description.strip():
                    job_description = f"Job position: {job_name}"
            else:
                # Fallback if job description file doesn't exist
                job_description = f"Job position: {job_name}"
                
            # Cache the job description
            self.job_description_cache[job_description_file] = job_description
        
        # Process text with BERT tokenizer - using job description and CV text
        encoding = self.bert_tokenizer(
            job_description,  # First text (Job description)
            cv_text,          # Second text (CV content)
            add_special_tokens=True,
            max_length=self.max_length,
            padding="max_length",
            truncation=True,
            return_tensors="pt",
            return_overflowing_tokens=False,  # Don't return overflowing tokens
        )
        
        # Get tokenized data
        input_ids = encoding["input_ids"].squeeze()
        attention_mask = encoding["attention_mask"].squeeze()
        token_type_ids = encoding["token_type_ids"].squeeze()
        
        # Create image from CV file
        image = create_page_image(cv_path)
        
        # Process image for ViT
        visual_inputs = self.vit_processor(images=image, return_tensors="pt")
        pixel_values = visual_inputs.pixel_values.squeeze()
        
        # One-hot encode job name
        job_idx = self.job_to_idx[job_name]
        job_one_hot = torch.zeros(len(self.unique_jobs), dtype=torch.float32)
        job_one_hot[job_idx] = 1.0
        
        # Format features tensor (default all zeros)
        format_features = torch.zeros(4, dtype=torch.float32)
        
        # Extract formatting features if enabled
        if self.extract_format_features:
            # Check cache first
            if cv_path in self.format_features_cache:
                # Agr Cv mili jiska format nikal chukay hai to wohi cv lay aoo 
                features = self.format_features_cache[cv_path]
            else:
                # CV kay format features lay kay aoo 
                features = extract_cv_features(cv_path)
                self.format_features_cache[cv_path] = features
                
            # Normalize and convert to tensor
            format_features = torch.tensor([
                min(features["font_count"], 10) / 10.0,   # Normalized font count
                min(features["text_density"], 20) / 20.0, # Normalized text density
                min(features["bullet_points"], 30) / 30.0, # Normalized bullet points
                min(features["pages"], 5) / 5.0           # Normalized page count
            ], dtype=torch.float32)
        
        return {
            # BERT inputs
            "input_ids": input_ids,
            "attention_mask": attention_mask,
            "token_type_ids": token_type_ids,
            # ViT inputs
            "pixel_values": pixel_values,
            # Additional features
            "format_features": format_features,
            "job_one_hot": job_one_hot,
            # Label
            "label": torch.tensor(label, dtype=torch.long),
            # Metadata (for reference)
            "cv_path": cv_path,
            "job_name": job_name,
            "status": status
        }


class MultimodalCVClassifier(nn.Module):
    """
    Advanced multimodal CV classifier with separate branches and attention-based fusion
    
    Architecture:
    1. Visual Analysis Branch:
       - ViT for document appearance features
       - FC layers for visual score
       
    2. Semantic Analysis Branch:
       - BERT for CV-JD text matching
       - FC layers for semantic score
       
    3. Format Features Branch:
       - Additional extracted document format features
       - FC layers for processing
       
    4. Job Information Branch:
       - One-hot encoded job information
       - FC layers for processing
       
    5. Multimodal Fusion:
       - Attention-weighted combination of branches
       - Classifier for final decision
    """
    def __init__(self, vit_model_name, bert_model_name, num_classes, num_jobs, dropout_rate=0.3):
        super(MultimodalCVClassifier, self).__init__()
        
        # Load pre-trained ViT model for visual analysis
        self.vit = ViTModel.from_pretrained(vit_model_name)
        
        # Load pre-trained BERT model for semantic analysis
        self.bert = BertModel.from_pretrained(bert_model_name)
        
        # Get feature dimensions
        self.vit_dim = self.vit.config.hidden_size
        self.bert_dim = self.bert.config.hidden_size
        self.format_dim = 4  # Number of format features
        self.job_dim = num_jobs  # Number of unique jobs
        
        # Visual analysis branch (document appearance)
        self.visual_branch = nn.Sequential(
            nn.Linear(self.vit_dim, 512),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(512, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Visual score prediction (how visually appealing is the CV)
        self.visual_score = nn.Linear(256, 1)
        
        # Semantic analysis branch (CV-JD matching)
        # Use larger hidden dim for better understanding
        self.semantic_branch = nn.Sequential(
            nn.Linear(self.bert_dim * 2, 768),  # CV + JD embeddings concatenated
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(768, 512),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(512, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Cross-attention between CV and JD
        self.cv_jd_attention = nn.MultiheadAttention(
            embed_dim=self.bert_dim,
            num_heads=8,
            dropout=dropout_rate
        )
        
        # Semantic score prediction (how well CV matches JD)
        self.semantic_score = nn.Linear(256, 1)
        
        # Format features processing
        self.format_branch = nn.Sequential(
            nn.Linear(self.format_dim, 32),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(32, 64),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Job information processing
        self.job_branch = nn.Sequential(
            nn.Linear(self.job_dim, 32),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(32, 64),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Combined model for final classification
        self.combined_dim = 256 + 256 + 64 + 64  # visual + semantic + format + job
        
        # Attention fusion for adaptive weighting of different branches
        self.attention = nn.Sequential(
            nn.Linear(self.combined_dim, 4),  # 4 attention weights for 4 branches
            nn.Softmax(dim=1)  # Normalize attention weights
        )
        
        # Final classification layer
        self.classifier = nn.Sequential(
            nn.Linear(self.combined_dim, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(256, num_classes)
        )
    
    def forward(self, pixel_values, input_ids, attention_mask, token_type_ids, 
                format_features=None, job_one_hot=None):
        """
        Forward pass through the multimodal network with attention fusion
        
        Args:
            pixel_values: Image tensor for ViT
            input_ids: Token IDs for BERT
            attention_mask: Attention mask for BERT
            token_type_ids: Token type IDs for BERT
            format_features: Additional CV formatting features
            job_one_hot: One-hot encoded job information
            
        Returns:
            dict: Dictionary with logits, visual score, semantic score, and attention weights
        """
        # Process format features (default to zeros if not provided)
        if format_features is None:
            format_features = torch.zeros(pixel_values.size(0), self.format_dim, 
                                         device=pixel_values.device)
        
        # Process job features (default to zeros if not provided)
        if job_one_hot is None:
            job_one_hot = torch.zeros(pixel_values.size(0), self.job_dim,
                                     device=pixel_values.device)
        
        # Visual branch: Extract visual features with ViT
        vit_outputs = self.vit(pixel_values=pixel_values)
        vit_features = vit_outputs.pooler_output  # [batch_size, vit_dim]
        visual_features = self.visual_branch(vit_features)  # [batch_size, 256]
        
        # Calculate visual score (higher = more visually appealing)
        visual_score_val = self.visual_score(visual_features)  # [batch_size, 1]
        
        # Semantic branch: Extract text features with BERT
        bert_outputs = self.bert(
            input_ids=input_ids,
            attention_mask=attention_mask,
            token_type_ids=token_type_ids,
        )
        
        # Get all hidden states for cross-attention
        all_hidden_states = bert_outputs.last_hidden_state  # [batch_size, seq_len, bert_dim]
        
        # Separate CV and JD representations using token_type_ids
        # token_type_ids = 0 for JD, 1 for CV
        jd_mask = (token_type_ids == 0).float()
        cv_mask = (token_type_ids == 1).float()
        
        # Get average embeddings for JD and CV
        jd_embeddings = (all_hidden_states * jd_mask.unsqueeze(-1)).sum(dim=1) / jd_mask.sum(dim=1, keepdim=True).clamp(min=1)
        cv_embeddings = (all_hidden_states * cv_mask.unsqueeze(-1)).sum(dim=1) / cv_mask.sum(dim=1, keepdim=True).clamp(min=1)
        
        # Apply cross-attention between CV and JD
        cv_attended, _ = self.cv_jd_attention(
            cv_embeddings.unsqueeze(0),  # Query
            jd_embeddings.unsqueeze(0),  # Key
            jd_embeddings.unsqueeze(0),  # Value
        )
        cv_attended = cv_attended.squeeze(0)
        
        # Concatenate original and attended features
        combined_semantic = torch.cat([cv_embeddings, cv_attended], dim=1)  # [batch_size, bert_dim * 2]
        
        semantic_features = self.semantic_branch(combined_semantic)  # [batch_size, 256]
        
        # Calculate semantic score (higher = better match between CV and JD)
        semantic_score_val = self.semantic_score(semantic_features)  # [batch_size, 1]
        
        # Format branch: Process formatting features
        format_features = self.format_branch(format_features)  # [batch_size, 64]
        
        # Job branch: Process job information
        job_features = self.job_branch(job_one_hot)  # [batch_size, 64]
        
        # Concatenate all features for fusion
        combined_features = torch.cat(
            (visual_features, semantic_features, format_features, job_features), dim=1
        )  # [batch_size, combined_dim]
        
        # Calculate attention weights for adaptive fusion
        attention_weights = self.attention(combined_features)  # [batch_size, 4]
        
        # Human brain-like processing: 
        # If semantic score is high (good CV-JD match), give it more weight
        # If visual score is low (poor formatting), reduce its importance
        semantic_importance = torch.sigmoid(semantic_score_val).squeeze()
        visual_penalty = torch.sigmoid(-visual_score_val).squeeze() * 0.5
        
        # Adjust weights based on scores (like human decision making)
        adjusted_weights = attention_weights.clone()
        adjusted_weights[:, 1] = adjusted_weights[:, 1] + semantic_importance * 0.3  # Boost semantic
        adjusted_weights[:, 0] = adjusted_weights[:, 0] - visual_penalty * 0.2  # Reduce visual if poor
        
        # Re-normalize to sum to 1
        attention_weights = F.softmax(adjusted_weights, dim=1)
        
        # Apply attention weights (reshape for broadcasting)
        visual_weight = attention_weights[:, 0].unsqueeze(1).expand_as(visual_features)
        semantic_weight = attention_weights[:, 1].unsqueeze(1).expand_as(semantic_features)
        format_weight = attention_weights[:, 2].unsqueeze(1).expand_as(format_features)
        job_weight = attention_weights[:, 3].unsqueeze(1).expand_as(job_features)
        
        # Weight each modality
        weighted_visual = visual_features * visual_weight
        weighted_semantic = semantic_features * semantic_weight
        weighted_format = format_features * format_weight
        weighted_job = job_features * job_weight
        
        # Concatenate weighted features
        weighted_features = torch.cat(
            (weighted_visual, weighted_semantic, weighted_format, weighted_job), dim=1
        )
        
        # Final classification
        logits = self.classifier(weighted_features)  # [batch_size, num_classes]
        
        return {
            'logits': logits,
            'visual_score': visual_score_val,
            'semantic_score': semantic_score_val,
            'attention_weights': attention_weights
        }


def train_model(model, train_dataloader, val_dataloader, optimizer, criterion, epochs, device):
    """
    Train the multimodal model with visual and semantic loss components
    
    This extended training includes:
    - Classification loss: CrossEntropy for primary task
    - Visual-semantic consistency loss: To ensure correlation between branches
    - Regularization: To prevent overfitting
    
    Args:
        model: Model to train
        train_dataloader: DataLoader for training data
        val_dataloader: DataLoader for validation data
        optimizer: Optimizer for training
        criterion: Main classification loss function
        epochs (int): Number of training epochs
        device: Device to train on (CPU/GPU)
        
    Returns:
        MultimodalCVClassifier: Trained model
    """
    model.to(device)
    best_val_accuracy = 0.0
    
    print(f"Training on device: {device}")
    print(f"Starting training for {epochs} epochs")

    for epoch in range(epochs):
        # Training phase
        model.train()
        train_loss = 0.0
        classification_loss_total = 0.0
        visual_loss_total = 0.0
        semantic_loss_total = 0.0
        attention_loss_total = 0.0
        correct_predictions = 0
        total_predictions = 0

        print(f"\nEpoch {epoch+1}/{epochs}")
        print("Training...")
        
        for batch_idx, batch in enumerate(train_dataloader):
            # Move data to device
            input_ids = batch["input_ids"].to(device)
            attention_mask = batch["attention_mask"].to(device)
            token_type_ids = batch["token_type_ids"].to(device)
            pixel_values = batch["pixel_values"].to(device)
            format_features = batch["format_features"].to(device)
            job_one_hot = batch["job_one_hot"].to(device)
            labels = batch["label"].to(device)

            # Clear gradients
            optimizer.zero_grad()

            # Forward pass
            outputs = model(
                pixel_values, 
                input_ids, 
                attention_mask, 
                token_type_ids,
                format_features,
                job_one_hot
            )
            
            # Get logits and scores
            logits = outputs['logits']
            visual_score = outputs['visual_score'].squeeze()
            semantic_score = outputs['semantic_score'].squeeze()
            attention_weights = outputs['attention_weights']
            
            # Classification loss (main objective)
            classification_loss = criterion(logits, labels)
            
            # Visual consistency loss (better visuals = better classes)
            # Normalize labels to [0,1] range for regression loss
            normalized_labels = 1.0 - (labels.float() / 3.0)  # 0=1.0, 3=0.0
            visual_loss = F.mse_loss(torch.sigmoid(visual_score), normalized_labels)
            
            # Semantic consistency loss (better matches = better classes)
            semantic_loss = F.mse_loss(torch.sigmoid(semantic_score), normalized_labels)
            
            # Attention regularization to encourage semantic attention
            # Target: visual=0.2, semantic=0.6, format=0.1, job=0.1
            target_attention = torch.tensor([0.2, 0.6, 0.1, 0.1], device=device)
            target_attention = target_attention.unsqueeze(0).expand_as(attention_weights)
            attention_reg_loss = F.mse_loss(attention_weights, target_attention)
            
            # Combined loss with weighting
            loss = classification_loss + 0.2 * visual_loss + 0.3 * semantic_loss + 0.5 * attention_reg_loss
            
            # Backward pass and optimize
            loss.backward()
            optimizer.step()

            # Accumulate statistics
            train_loss += loss.item()
            classification_loss_total += classification_loss.item()
            visual_loss_total += visual_loss.item()
            semantic_loss_total += semantic_loss.item()
            attention_loss_total += attention_reg_loss.item()
            
            _, preds = torch.max(logits, dim=1)
            correct_predictions += torch.sum(preds == labels).item()
            total_predictions += labels.size(0)
            
            # Print batch progress every 5 batches
            if (batch_idx + 1) % 5 == 0:
                print(f"  Batch {batch_idx + 1}/{len(train_dataloader)}, "
                      f"Loss: {loss.item():.4f}, "
                      f"Accuracy: {torch.sum(preds == labels).item() / labels.size(0):.4f}")

        # Calculate training metrics
        train_accuracy = correct_predictions / total_predictions
        train_loss = train_loss / len(train_dataloader)
        classification_loss_avg = classification_loss_total / len(train_dataloader)
        visual_loss_avg = visual_loss_total / len(train_dataloader)
        semantic_loss_avg = semantic_loss_total / len(train_dataloader)
        attention_loss_avg = attention_loss_total / len(train_dataloader)

        # Validation phase
        val_accuracy, val_loss, val_component_losses = evaluate_model(
            model, val_dataloader, criterion, device
        )

        # Print epoch results
        print(f"Epoch {epoch+1}/{epochs} Results:")
        print(f"  Train Loss: {train_loss:.4f}, Train Accuracy: {train_accuracy:.4f}")
        print(f"  Train Loss Components: Class={classification_loss_avg:.4f}, Visual={visual_loss_avg:.4f}, Semantic={semantic_loss_avg:.4f}, Attention={attention_loss_avg:.4f}")
        print(f"  Val Loss: {val_loss:.4f}, Val Accuracy: {val_accuracy:.4f}")
        print(f"  Val Loss Components: Class={val_component_losses['class']:.4f}, Visual={val_component_losses['visual']:.4f}, Semantic={val_component_losses['semantic']:.4f}")

        # Save best model
        if val_accuracy > best_val_accuracy:
            best_val_accuracy = val_accuracy
            torch.save(model.state_dict(), "best_model.pt")
            print(f"  New best model saved with validation accuracy: {val_accuracy:.4f}")

    print("\nTraining complete!")
    print(f"Best validation accuracy: {best_val_accuracy:.4f}")
    
    return model


def evaluate_model(model, dataloader, criterion, device):
    """
    Evaluate the model on validation data with detailed loss components
    
    Args:
        model: Model to evaluate
        dataloader: DataLoader for evaluation data
        criterion: Loss function
        device: Device to evaluate on
        
    Returns:
        tuple: (accuracy, average_loss, component_losses_dict)
    """
    model.eval()
    val_loss = 0.0
    classification_loss_total = 0.0
    visual_loss_total = 0.0
    semantic_loss_total = 0.0
    correct_predictions = 0
    total_predictions = 0

    print("Evaluating...")
    
    with torch.no_grad():
        for batch in dataloader:
            # Move data to device
            input_ids = batch["input_ids"].to(device)
            attention_mask = batch["attention_mask"].to(device)
            token_type_ids = batch["token_type_ids"].to(device)
            pixel_values = batch["pixel_values"].to(device)
            format_features = batch["format_features"].to(device)
            job_one_hot = batch["job_one_hot"].to(device)
            labels = batch["label"].to(device)

            # Forward pass
            outputs = model(
                pixel_values, 
                input_ids, 
                attention_mask, 
                token_type_ids,
                format_features,
                job_one_hot
            )
            
            # Get logits and scores
            logits = outputs['logits']
            visual_score = outputs['visual_score'].squeeze()
            semantic_score = outputs['semantic_score'].squeeze()
            
            # Classification loss
            classification_loss = criterion(logits, labels)
            
            # Visual and semantic consistency losses
            normalized_labels = 1.0 - (labels.float() / 3.0)  # 0=1.0, 3=0.0
            visual_loss = F.mse_loss(torch.sigmoid(visual_score), normalized_labels)
            semantic_loss = F.mse_loss(torch.sigmoid(semantic_score), normalized_labels)
            
            # Combined loss with weighting
            loss = classification_loss + 0.2 * visual_loss + 0.2 * semantic_loss

            # Accumulate statistics
            val_loss += loss.item()
            classification_loss_total += classification_loss.item()
            visual_loss_total += visual_loss.item()
            semantic_loss_total += semantic_loss.item()
            
            _, preds = torch.max(logits, dim=1)
            correct_predictions += torch.sum(preds == labels).item()
            total_predictions += labels.size(0)

    # Calculate metrics
    accuracy = correct_predictions / total_predictions if total_predictions > 0 else 0
    average_loss = val_loss / len(dataloader) if len(dataloader) > 0 else 0
    
    component_losses = {
        'class': classification_loss_total / len(dataloader) if len(dataloader) > 0 else 0,
        'visual': visual_loss_total / len(dataloader) if len(dataloader) > 0 else 0,
        'semantic': semantic_loss_total / len(dataloader) if len(dataloader) > 0 else 0
    }

    return accuracy, average_loss, component_losses


def predict_cv(model, cv_path, job_name, job_description_file, job_one_hot, bert_tokenizer, vit_processor, device):
    """
    Make prediction for a single CV with detailed scores
    
    Args:
        model: Trained model
        cv_path: Path to CV file
        job_name: Name of the job
        job_description_file: Path to the job description file
        job_one_hot: One-hot encoded job vector
        bert_tokenizer: BERT tokenizer
        vit_processor: ViT image processor
        device: Device to run inference on
        
    Returns:
        dict: Prediction results including class, probabilities and component scores
    """
    model.eval()
    
    # Create image from CV file
    image = create_page_image(cv_path)
    
    # Extract CV text
    cv_text = extract_text_from_file(cv_path)
    
    # Read the actual job description from file
    if job_description_file and os.path.exists(job_description_file):
        job_description = extract_text_from_file(job_description_file)
        # If job description is empty, use fallback
        if not job_description.strip():
            job_description = f"Job position: {job_name}"
    else:
        # Fallback if job description file doesn't exist
        job_description = f"Job position: {job_name}"
    
    # Extract format features
    format_features_dict = extract_cv_features(cv_path)
    format_features = torch.tensor([
        min(format_features_dict["font_count"], 10) / 10.0,
        min(format_features_dict["text_density"], 20) / 20.0,
        min(format_features_dict["bullet_points"], 30) / 30.0,
        min(format_features_dict["pages"], 5) / 5.0
    ], dtype=torch.float32).unsqueeze(0).to(device)
    
    # Process text with BERT tokenizer
    encoding = bert_tokenizer(
        job_description,
        cv_text,
        add_special_tokens=True,
        max_length=256,
        padding="max_length",
        truncation=True,
        return_tensors="pt",
        return_overflowing_tokens=False,  # Don't return overflowing tokens
    )
    
    # Process image for ViT
    visual_inputs = vit_processor(images=image, return_tensors="pt")
    
    # Move inputs to device
    input_ids = encoding["input_ids"].to(device)
    attention_mask = encoding["attention_mask"].to(device)
    token_type_ids = encoding["token_type_ids"].to(device)
    pixel_values = visual_inputs.pixel_values.to(device)
    job_one_hot = job_one_hot.unsqueeze(0).to(device)
    
    # Make prediction
    with torch.no_grad():
        outputs = model(
            pixel_values,
            input_ids,
            attention_mask,
            token_type_ids,
            format_features,
            job_one_hot
        )
        
        logits = outputs['logits']
        visual_score = outputs['visual_score'].item()
        semantic_score = outputs['semantic_score'].item()
        attention_weights = outputs['attention_weights'][0].cpu().numpy()
        
        # Get probabilities and predicted class
        probabilities = F.softmax(logits, dim=1)[0]
        pred_class = torch.argmax(probabilities).item()
        
        # Map back to status
        status_map_reverse = {
            0: "ACCEPT",
            1: "INTERVIEW",
            2: "SHORTLIST", 
            3: "REJECT",
        }
        
        status = status_map_reverse.get(pred_class, "Unknown")
        
    # Return detailed results
    result = {
        "status": status,
        "class_id": pred_class,
        "confidence": probabilities[pred_class].item(),
        "probabilities": {status_map_reverse.get(i, "Unknown"): prob.item() 
                         for i, prob in enumerate(probabilities)},
        "visual_score": sigmoid(visual_score),  # Normalized to [0,1]
        "semantic_score": sigmoid(semantic_score),  # Normalized to [0,1]
        "attention_weights": {
            "visual": attention_weights[0],
            "semantic": attention_weights[1],
            "format": attention_weights[2],
            "job": attention_weights[3]
        }
    }
    
    return result


def sigmoid(x):
    """
    Sigmoid function to normalize scores to [0,1] range
    """
    return 1 / (1 + np.exp(-x))


def main():
    """
    Main function to train the multimodal CV classification model
    """
    print("Initializing multimodal CV classification training...")
    print(f"PyMuPDF version: {fitz.version[0]}")
    print(f"PIL version: {Image.__version__ if hasattr(Image, '__version__') else 'Unknown'}")
    
    # Set random seed for reproducibility
    torch.manual_seed(42)
    np.random.seed(42)
    
    # Set device
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")

    # Load tokenizer and image processor
    print("Loading tokenizers and processors...")
    bert_tokenizer = BertTokenizer.from_pretrained(BERT_MODEL_PATH)
    vit_processor = ViTImageProcessor.from_pretrained(VIT_MODEL_PATH)

    # Load CV data from folder structure
    try:
        print(f"\nLoading CV data from: {CV_FILES_DIR}")
        cv_data = load_cv_data_from_folder_structure(CV_FILES_DIR)
        print(f"\nSuccessfully loaded {len(cv_data)} CV files for analysis")
        
        # Group by job and status for better overview
        job_status_counts = {}
        job_totals = {}
        for _, job, status, _, _ in cv_data:
            key = (job, status)
            job_status_counts[key] = job_status_counts.get(key, 0) + 1
            job_totals[job] = job_totals.get(job, 0) + 1
        
        print("\nDataset distribution:")
        for job in sorted(job_totals.keys()):
            print(f"\n  {job}: {job_totals[job]} total CVs")
            for status in ['ACCEPT', 'INTERVIEW', 'SHORTLIST', 'REJECT']:
                count = job_status_counts.get((job, status), 0)
                if count > 0:
                    print(f"    - {status}: {count} CVs")
            
    except Exception as e:
        print(f"\nError loading CV data: {str(e)}")
        import traceback
        traceback.print_exc()
        return

    # Create dataset
    print("\nCreating dataset with CV files...")
    print("Using enhanced configuration:")
    print("- Max token length: 512 (increased from 256)")
    print("- Chunk processing: Enabled (for long documents)")
    print("- Attention regularization: Targeting 60% semantic weight")
    dataset = CVJobDataset(
        cv_data=cv_data,
        vit_processor=vit_processor,
        bert_tokenizer=bert_tokenizer,
        max_length=512,
        extract_format_features=True,
        use_chunk_processing=True
    )

    # Split into train and validation (with more data in training)
    dataset_size = len(dataset)
    train_size = int(0.8 * dataset_size)
    val_size = dataset_size - train_size
    
    # Make sure we have at least one sample in each set
    if train_size == 0:
        train_size = 1
        val_size = max(0, dataset_size - 1)
    
    if val_size == 0 and dataset_size > 1:
        val_size = 1
        train_size = dataset_size - 1
    
    # Use random_split with generator for reproducibility
    generator = torch.Generator().manual_seed(42)
    if dataset_size > 1:
        train_dataset, val_dataset = random_split(
            dataset, [train_size, val_size], generator=generator
        )
    else:
        # If we only have one sample, use it for both training and validation
        train_dataset = dataset
        val_dataset = dataset

    print(f"\nSplit dataset: {train_size} for training, {val_size} for validation")

    # Create data loaders (with smaller batch size due to limited data)
    batch_size = min(32, train_size)  # Make sure batch size is not larger than dataset
    train_dataloader = DataLoader(train_dataset, batch_size=batch_size, shuffle=True)
    val_dataloader = DataLoader(val_dataset, batch_size=batch_size)

    # Initialize model
    print("\nInitializing multimodal model...")
    num_classes = 4  # Four possible statuses
    num_jobs = len(dataset.unique_jobs)
    model = MultimodalCVClassifier(VIT_MODEL_PATH, BERT_MODEL_PATH, num_classes, num_jobs)

    # Print model architecture summary
    print("\nModel Architecture:")
    print(f"- Visual Branch: ViT → {model.vit_dim} → 512 → 256 → 1")
    print(f"- Semantic Branch: BERT → {model.bert_dim} → 512 → 256 → 1")
    print(f"- Format Branch: {model.format_dim} → 32 → 64")
    print(f"- Job Branch: {model.job_dim} → 32 → 64")
    print(f"- Fusion: Attention-weighted combination → 256 → {num_classes}")
    print(f"- Total Parameters: {sum(p.numel() for p in model.parameters() if p.requires_grad):,}")

    # Calculate class weights for imbalanced data
    print("\nCalculating class weights for balanced training...")
    class_counts = torch.zeros(4)  # 4 classes: ACCEPT, INTERVIEW, SHORTLIST, REJECT
    for i in range(len(train_dataset)):
        label = train_dataset[i]['label']
        class_counts[label] += 1
    
    # Calculate inverse frequency weights
    total_samples = class_counts.sum()
    class_weights = total_samples / (4 * class_counts)
    class_weights = class_weights / class_weights.mean()  # Normalize
    
    print("Class weights:")
    status_names = ["ACCEPT", "INTERVIEW", "SHORTLIST", "REJECT"]
    for i, (name, weight) in enumerate(zip(status_names, class_weights)):
        print(f"  {name}: {weight:.2f} (count: {int(class_counts[i])})")
    
    # Define optimizer and loss function (with weight decay for regularization)
    optimizer = optim.AdamW(model.parameters(), lr=2e-5, weight_decay=0.01)
    criterion = nn.CrossEntropyLoss(weight=class_weights.to(device))

    # Determine number of epochs (fewer if we have very limited data)
    if dataset_size <= 2:
        epochs = 50  # More epochs for very small datasets
    elif dataset_size <= 10:
        epochs = 30  # More epochs for small datasets
    else:
        epochs = 20  # Fewer epochs for larger datasets

    # Train the model
    print(f"\nStarting training for {epochs} epochs...")
    model = train_model(
        model, 
        train_dataloader, 
        val_dataloader, 
        optimizer, 
        criterion, 
        epochs=epochs,
        device=device
    )

    # Save the final model
    torch.save(model.state_dict(), 'final_model.pt')
    print("Final model saved to final_model.pt")
    
    # Test the model on each CV file
    print("\nTesting model on all CV files:")
    for i, (cv_path, job_name, status, _, job_description_file) in enumerate(cv_data):
        # Get job one-hot encoding
        job_idx = dataset.job_to_idx[job_name]
        job_one_hot = torch.zeros(len(dataset.unique_jobs), dtype=torch.float32)
        job_one_hot[job_idx] = 1.0
        
        # Predict
        result = predict_cv(
            model, cv_path, job_name, job_description_file, job_one_hot,
            bert_tokenizer, vit_processor, device
        )
        
        # Print results
        print(f"\n{i+1}. CV: {os.path.basename(cv_path)}")
        print(f"   Job: {job_name}, True Status: {status}")
        print(f"   Predicted Status: {result['status']} (Confidence: {result['confidence']:.2%})")
        print(f"   Visual Score: {result['visual_score']:.2%}")
        print(f"   Semantic Score: {result['semantic_score']:.2%}")
        print(f"   Attention Weights: Visual={result['attention_weights']['visual']:.2f}, " 
              f"Semantic={result['attention_weights']['semantic']:.2f}, "
              f"Format={result['attention_weights']['format']:.2f}, "
              f"Job={result['attention_weights']['job']:.2f}")


if __name__ == "__main__":
    main()