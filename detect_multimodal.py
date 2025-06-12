import torch
from torch import nn
import torch.nn.functional as F
from transformers import BertModel, BertTokenizer, ViTModel, ViTImageProcessor
import os
import sys
import argparse
from pathlib import Path
import numpy as np
from PIL import Image
import fitz  # PyMuPDF
import io
import re

# Force offline mode
os.environ["TRANSFORMERS_OFFLINE"] = "1"

# Copy the exact model architecture from train.py
class MultimodalCVClassifier(nn.Module):
    def __init__(self, vit_model_name, bert_model_name, num_classes, num_jobs, dropout_rate=0.3):
        super(MultimodalCVClassifier, self).__init__()
        
        # Load pre-trained ViT model for visual analysis
        self.vit = ViTModel.from_pretrained(vit_model_name)
        
        # Load pre-trained BERT model for semantic analysis
        self.bert = BertModel.from_pretrained(bert_model_name)
        
        # Get feature dimensions
        self.vit_dim = self.vit.config.hidden_size
        self.bert_dim = self.bert.config.hidden_size
        self.format_dim = 4  # Number of format features
        self.job_dim = num_jobs  # Number of unique jobs
        
        # Visual analysis branch (document appearance)
        self.visual_branch = nn.Sequential(
            nn.Linear(self.vit_dim, 512),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(512, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Visual score prediction
        self.visual_score = nn.Linear(256, 1)
        
        # Semantic analysis branch with cross-attention
        self.semantic_branch = nn.Sequential(
            nn.Linear(self.bert_dim * 2, 768),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(768, 512),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(512, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Cross-attention between CV and JD
        self.cv_jd_attention = nn.MultiheadAttention(
            embed_dim=self.bert_dim,
            num_heads=8,
            dropout=dropout_rate
        )
        
        # Semantic score prediction
        self.semantic_score = nn.Linear(256, 1)
        
        # Format features processing
        self.format_branch = nn.Sequential(
            nn.Linear(self.format_dim, 32),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(32, 64),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Job information processing
        self.job_branch = nn.Sequential(
            nn.Linear(self.job_dim, 32),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(32, 64),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
        )
        
        # Combined model for final classification
        self.combined_dim = 256 + 256 + 64 + 64
        
        # Attention fusion
        self.attention = nn.Sequential(
            nn.Linear(self.combined_dim, 4),
            nn.Softmax(dim=1)
        )
        
        # Final classification layer
        self.classifier = nn.Sequential(
            nn.Linear(self.combined_dim, 256),
            nn.ReLU(),
            nn.Dropout(dropout_rate),
            nn.Linear(256, num_classes),
        )

    def forward(self, pixel_values, input_ids, attention_mask, token_type_ids=None, 
                format_features=None, job_one_hot=None):
        
        # Default values
        if format_features is None:
            format_features = torch.zeros(pixel_values.size(0), self.format_dim, 
                                         device=pixel_values.device)
        
        if job_one_hot is None:
            job_one_hot = torch.zeros(pixel_values.size(0), self.job_dim,
                                     device=pixel_values.device)
        
        # Visual branch
        vit_outputs = self.vit(pixel_values=pixel_values)
        vit_features = vit_outputs.pooler_output
        visual_features = self.visual_branch(vit_features)
        visual_score_val = self.visual_score(visual_features)
        
        # Semantic branch with cross-attention
        bert_outputs = self.bert(
            input_ids=input_ids,
            attention_mask=attention_mask,
            token_type_ids=token_type_ids,
        )
        
        all_hidden_states = bert_outputs.last_hidden_state
        
        # Separate CV and JD representations
        jd_mask = (token_type_ids == 0).float()
        cv_mask = (token_type_ids == 1).float()
        
        jd_embeddings = (all_hidden_states * jd_mask.unsqueeze(-1)).sum(dim=1) / jd_mask.sum(dim=1, keepdim=True).clamp(min=1)
        cv_embeddings = (all_hidden_states * cv_mask.unsqueeze(-1)).sum(dim=1) / cv_mask.sum(dim=1, keepdim=True).clamp(min=1)
        
        # Cross-attention
        cv_attended, _ = self.cv_jd_attention(
            cv_embeddings.unsqueeze(0),
            jd_embeddings.unsqueeze(0),
            jd_embeddings.unsqueeze(0),
        )
        cv_attended = cv_attended.squeeze(0)
        
        combined_semantic = torch.cat([cv_embeddings, cv_attended], dim=1)
        semantic_features = self.semantic_branch(combined_semantic)
        semantic_score_val = self.semantic_score(semantic_features)
        
        # Format and job branches
        format_features = self.format_branch(format_features)
        job_features = self.job_branch(job_one_hot)
        
        # Concatenate all features
        combined_features = torch.cat(
            (visual_features, semantic_features, format_features, job_features), dim=1
        )
        
        # Calculate attention weights
        attention_weights = self.attention(combined_features)
        
        # Human brain-like processing
        semantic_importance = torch.sigmoid(semantic_score_val).squeeze()
        visual_penalty = torch.sigmoid(-visual_score_val).squeeze() * 0.5
        
        adjusted_weights = attention_weights.clone()
        adjusted_weights[:, 1] = adjusted_weights[:, 1] + semantic_importance * 0.3
        adjusted_weights[:, 0] = adjusted_weights[:, 0] - visual_penalty * 0.2
        
        attention_weights = F.softmax(adjusted_weights, dim=1)
        
        # Apply attention weights
        visual_weight = attention_weights[:, 0].unsqueeze(1).expand_as(visual_features)
        semantic_weight = attention_weights[:, 1].unsqueeze(1).expand_as(semantic_features)
        format_weight = attention_weights[:, 2].unsqueeze(1).expand_as(format_features)
        job_weight = attention_weights[:, 3].unsqueeze(1).expand_as(job_features)
        
        # Weight each modality
        weighted_visual = visual_features * visual_weight
        weighted_semantic = semantic_features * semantic_weight
        weighted_format = format_features * format_weight
        weighted_job = job_features * job_weight
        
        # Concatenate weighted features
        weighted_features = torch.cat(
            (weighted_visual, weighted_semantic, weighted_format, weighted_job), dim=1
        )
        
        # Final classification
        logits = self.classifier(weighted_features)
        
        return {
            'logits': logits,
            'visual_score': visual_score_val,
            'semantic_score': semantic_score_val,
            'attention_weights': attention_weights
        }


def create_page_image(file_path, page_num=0):
    """Create image from first page of document"""
    if file_path.lower().endswith('.pdf'):
        try:
            doc = fitz.open(file_path)
            if len(doc) > 0:
                page = doc[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                img_data = pix.pil_tobytes("RGB")
                img = Image.open(io.BytesIO(img_data))
                doc.close()
                return img
            doc.close()
        except:
            pass
    
    # Return blank image for non-PDFs or errors
    return Image.new('RGB', (224, 224), color='white')


def extract_text_from_file(file_path):
    """Extract text from various file formats"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        
        # PDF files
        if ext == '.pdf':
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
            
        # Text files
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
                
        # Word documents
        elif ext in ['.docx', '.doc']:
            try:
                import docx
                doc = docx.Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            except:
                return f"Error reading Word document: {file_path}"
                
        # PowerPoint presentations
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            except:
                return f"Failed to read PowerPoint: {file_path}"
                
        # OpenDocument Text
        elif ext == '.odt':
            try:
                from odf import text, teletype
                from odf.opendocument import load
                doc = load(file_path)
                allparas = doc.getElementsByType(text.P)
                return "\n".join([teletype.extractText(p) for p in allparas])
            except:
                return f"Failed to read ODT: {file_path}"
                
        # Email files
        elif ext == '.eml':
            try:
                import email
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    msg = email.message_from_file(f)
                    body = ""
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            body += part.get_payload(decode=True).decode('utf-8', errors='replace')
                    return body if body else f"No text content in email: {file_path}"
            except:
                return f"Failed to read email: {file_path}"
                
        # MHTML files
        elif ext == '.mht':
            try:
                import re
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    content = f.read()
                    # Basic extraction - remove HTML tags
                    text = re.sub('<.*?>', '', content)
                    return text
            except:
                return f"Failed to read MHT: {file_path}"
                
        # HEIC images - skip text extraction
        elif ext == '.heic':
            return f"Image file (no text): {file_path}"
            
        # Binary files
        elif ext == '.bin':
            return f"Binary file (no text): {file_path}"
            
        # Default fallback
        else:
            return f"File type {ext} - attempting basic text read"
            
    except Exception as e:
        return f"Error reading file: {str(e)}"


def extract_cv_features(file_path):
    """Extract formatting features (simplified)"""
    # Return default features for now
    return {
        "font_count": 2,
        "text_density": 10,
        "bullet_points": 5,
        "pages": 1
    }


def evaluate_cv(jd_path, cv_path, model, tokenizer, vit_processor, device):
    """Evaluate CV against job description using multimodal model"""
    
    # Status mapping
    status_labels = {
        0: "ACCEPT",
        1: "INTERVIEW",
        2: "SHORTLIST",
        3: "REJECT",
    }
    
    # Extract texts
    jd_text = extract_text_from_file(jd_path)
    cv_text = extract_text_from_file(cv_path)
    
    # Create CV image
    cv_image = create_page_image(cv_path)
    
    # Process text with BERT tokenizer
    encoding = tokenizer(
        jd_text,
        cv_text,
        add_special_tokens=True,
        max_length=768,
        padding="max_length",
        truncation=True,
        return_tensors="pt",
    )
    
    # Process image with ViT processor
    visual_inputs = vit_processor(images=cv_image, return_tensors="pt")
    
    # Extract format features
    features = extract_cv_features(cv_path)
    format_features = torch.tensor([
        min(features["font_count"], 10) / 10.0,
        min(features["text_density"], 20) / 20.0,
        min(features["bullet_points"], 30) / 30.0,
        min(features["pages"], 5) / 5.0
    ], dtype=torch.float32).unsqueeze(0)
    
    # Move to device
    input_ids = encoding["input_ids"].to(device)
    attention_mask = encoding["attention_mask"].to(device)
    token_type_ids = encoding["token_type_ids"].to(device)
    pixel_values = visual_inputs["pixel_values"].to(device)
    format_features = format_features.to(device)
    
    # Evaluate
    model.eval()
    with torch.no_grad():
        outputs = model(
            pixel_values=pixel_values,
            input_ids=input_ids,
            attention_mask=attention_mask,
            token_type_ids=token_type_ids,
            format_features=format_features
        )
        
        logits = outputs['logits']
        visual_score = outputs['visual_score'].item()
        semantic_score = outputs['semantic_score'].item()
        attention_weights = outputs['attention_weights'][0]
        
        # Get prediction
        probabilities = F.softmax(logits, dim=1)
        predicted_class = torch.argmax(probabilities, dim=1).item()
        confidence = probabilities[0][predicted_class].item()
    
    return {
        "status": status_labels[predicted_class],
        "confidence": confidence,
        "visual_score": 1 / (1 + np.exp(-visual_score)),  # Sigmoid
        "semantic_score": 1 / (1 + np.exp(-semantic_score)),  # Sigmoid
        "probabilities": {
            status_labels[i]: prob.item() 
            for i, prob in enumerate(probabilities[0])
        },
        "attention_weights": {
            "visual": attention_weights[0].item(),
            "semantic": attention_weights[1].item(),
            "format": attention_weights[2].item(),
            "job": attention_weights[3].item()
        }
    }


def main():
    parser = argparse.ArgumentParser(description='Evaluate CV against Job Description')
    parser.add_argument('--jd', type=str, required=True, help='Path to job description file')
    parser.add_argument('--cv', type=str, required=True, help='Path to CV file')
    parser.add_argument('--model', type=str, default='best_model.pt', help='Path to model file')
    
    args = parser.parse_args()
    
    # Setup
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")
    
    # Load tokenizers and processors
    bert_model_path = "./local_bert_model"
    vit_model_path = "./local_vit_model"
    
    print("Loading tokenizers and processors...")
    tokenizer = BertTokenizer.from_pretrained(bert_model_path)
    vit_processor = ViTImageProcessor.from_pretrained(vit_model_path)
    
    # Initialize model
    num_classes = 4
    num_jobs = 10  # Based on actual training data
    model = MultimodalCVClassifier(
        vit_model_name=vit_model_path,
        bert_model_name=bert_model_path,
        num_classes=num_classes,
        num_jobs=num_jobs
    )
    
    # Load weights
    if os.path.exists(args.model):
        model.load_state_dict(torch.load(args.model, map_location=device))
        print(f"Model loaded from {args.model}")
    else:
        print(f"Warning: Model file {args.model} not found!")
        return
    
    model.to(device)
    
    # Evaluate
    print(f"\nEvaluating CV: {args.cv}")
    print(f"Against JD: {args.jd}")
    
    result = evaluate_cv(args.jd, args.cv, model, tokenizer, vit_processor, device)
    
    # Display results
    print("\n" + "="*60)
    print("EVALUATION RESULTS")
    print("="*60)
    print(f"Decision: {result['status']}")
    print(f"Confidence: {result['confidence']:.2%}")
    print(f"\nScores:")
    print(f"  Visual Score: {result['visual_score']:.2%}")
    print(f"  Semantic Score: {result['semantic_score']:.2%}")
    print(f"\nAttention Weights:")
    for key, weight in result['attention_weights'].items():
        print(f"  {key.capitalize()}: {weight:.2%}")
    print(f"\nProbabilities:")
    for status, prob in result['probabilities'].items():
        print(f"  {status}: {prob:.2%}")
    print("="*60)


if __name__ == "__main__":
    if len(sys.argv) == 1:
        print("Usage: python detect_multimodal.py --jd <job_file> --cv <cv_file>")
        print("Example: python detect_multimodal.py --jd job.txt --cv resume.pdf")
        sys.exit(1)
    
    main()